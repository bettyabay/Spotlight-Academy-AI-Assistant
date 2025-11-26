"""
Document processing for various file formats
"""
import os
from pathlib import Path
from typing import List, Dict
import logging

# PDF processing
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# DOCX processing
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# PPTX processing
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Image OCR
try:
    from PIL import Image
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

from config import Config
from .chunking import TextChunker

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Processes various document formats and extracts text"""
    
    def __init__(self):
        self.chunker = TextChunker()
        self.supported_formats = Config.SUPPORTED_FORMATS
    
    def process_file(self, file_path: str, metadata: Dict = None) -> List[Dict]:
        """
        Process a file and return chunks with metadata
        
        Args:
            file_path: Path to the file
            metadata: Additional metadata (module, chapter, lesson, concept)
            
        Returns:
            List of chunk dictionaries
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_ext = file_path.suffix.lower()
        
        if file_ext not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {file_ext}")
        
        # Extract text based on file type
        if file_ext == ".pdf":
            text = self._extract_pdf(file_path)
        elif file_ext == ".docx":
            text = self._extract_docx(file_path)
        elif file_ext == ".pptx":
            text = self._extract_pptx(file_path)
        elif file_ext in [".png", ".jpg", ".jpeg"]:
            text = self._extract_image_ocr(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_ext}")
        
        # Prepare metadata
        file_metadata = {
            "source_file": str(file_path.name),
            "file_path": str(file_path),
            "file_type": file_ext,
            **(metadata or {})
        }
        
        # Chunk the text
        chunks = self.chunker.chunk_text(text, file_metadata)
        
        logger.info(f"Processed {file_path.name}: {len(chunks)} chunks created")
        return chunks
    
    def _extract_pdf(self, file_path: Path) -> str:
        """Extract text from PDF file"""
        if not PDF_AVAILABLE:
            raise ImportError("PyPDF2 is required for PDF processing")
        
        text_content = []
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    if text.strip():
                        text_content.append(text)
        except Exception as e:
            logger.error(f"Error extracting PDF {file_path}: {str(e)}")
            raise
        
        return "\n\n".join(text_content)
    
    def _extract_docx(self, file_path: Path) -> str:
        """Extract text from DOCX file"""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx is required for DOCX processing")
        
        try:
            doc = Document(file_path)
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            return "\n\n".join(paragraphs)
        except Exception as e:
            logger.error(f"Error extracting DOCX {file_path}: {str(e)}")
            raise
    
    def _extract_pptx(self, file_path: Path) -> str:
        """Extract text from PPTX file"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is required for PPTX processing")
        
        try:
            prs = Presentation(file_path)
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if slide_text:
                    text_content.append(f"Slide {slide_num + 1}:\n" + "\n".join(slide_text))
            
            return "\n\n".join(text_content)
        except Exception as e:
            logger.error(f"Error extracting PPTX {file_path}: {str(e)}")
            raise
    
    def _extract_image_ocr(self, file_path: Path) -> str:
        """Extract text from image using OCR"""
        if not OCR_AVAILABLE:
            raise ImportError("Pillow and pytesseract are required for OCR")
        
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            return text
        except Exception as e:
            logger.error(f"Error extracting text from image {file_path}: {str(e)}")
            raise

